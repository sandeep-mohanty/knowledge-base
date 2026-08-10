#!/usr/bin/env python3
"""
Generate Passkey Authentication POC Presentation
Creates a professional PPTX for the CIAM passkey demo
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color scheme
COLORS = {
    'primary': RGBColor(0x14, 0x5a, 0x32),      # Dark green
    'secondary': RGBColor(0x1a, 0x25, 0x2f),    # Dark blue-gray
    'accent': RGBColor(0x27, 0xae, 0x60),       # Bright green
    'warning': RGBColor(0xe7, 0x4c, 0x3c),      # Red
    'text': RGBColor(0x2c, 0x3e, 0x50),         # Dark text
    'white': RGBColor(0xff, 0xff, 0xff),
    'light_bg': RGBColor(0xf5, 0xf5, 0xf5),
}

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Background shape
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = COLORS['primary']
    background.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(2.5), Inches(11.333), Inches(1.5)
    )
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS['white']
    title_para.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(
        Inches(1), Inches(4.2), Inches(11.333), Inches(1)
    )
    sub_frame = sub_box.text_frame
    sub_frame.text = subtitle
    sub_para = sub_frame.paragraphs[0]
    sub_para.font.size = Pt(24)
    sub_para.font.color.rgb = COLORS['accent']
    sub_para.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, content_items, two_column=False):
    """Add a content slide with title and bullet points"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Title bar
    title_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = COLORS['primary']
    title_bar.line.fill.background()
    
    # Title text
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(12), Inches(0.8)
    )
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS['white']
    
    # Content
    if two_column and len(content_items) >= 2:
        # Left column
        left_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.5), Inches(5.8), Inches(5.5)
        )
        left_frame = left_box.text_frame
        left_frame.word_wrap = True
        
        for i, item in enumerate(content_items[0]):
            if i == 0:
                p = left_frame.paragraphs[0]
            else:
                p = left_frame.add_paragraph()
            p.text = item
            p.font.size = Pt(18)
            p.font.color.rgb = COLORS['text']
            p.space_after = Pt(12)
            if item.startswith('✅') or item.startswith('❌') or item.startswith('⚠️'):
                p.level = 0
            else:
                p.level = 1
        
        # Right column
        right_box = slide.shapes.add_textbox(
            Inches(6.8), Inches(1.5), Inches(5.8), Inches(5.5)
        )
        right_frame = right_box.text_frame
        right_frame.word_wrap = True
        
        for i, item in enumerate(content_items[1]):
            if i == 0:
                p = right_frame.paragraphs[0]
            else:
                p = right_frame.add_paragraph()
            p.text = item
            p.font.size = Pt(18)
            p.font.color.rgb = COLORS['text']
            p.space_after = Pt(12)
    else:
        content_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.5), Inches(12), Inches(5.5)
        )
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        
        for i, item in enumerate(content_items):
            if i == 0:
                p = content_frame.paragraphs[0]
            else:
                p = content_frame.add_paragraph()
            p.text = item
            p.font.size = Pt(20)
            p.font.color.rgb = COLORS['text']
            p.space_after = Pt(14)
            if item.startswith('✅') or item.startswith('❌') or item.startswith('⚠️') or item.startswith('•'):
                p.level = 0
            else:
                p.level = 1
    
    return slide

def add_section_divider(prs, section_title, section_number):
    """Add a section divider slide"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = COLORS['secondary']
    background.line.fill.background()
    
    # Section number
    num_box = slide.shapes.add_textbox(
        Inches(1), Inches(2.8), Inches(11.333), Inches(1)
    )
    num_frame = num_box.text_frame
    num_frame.text = f"PART {section_number}"
    num_para = num_frame.paragraphs[0]
    num_para.font.size = Pt(32)
    num_para.font.bold = True
    num_para.font.color.rgb = COLORS['accent']
    num_para.alignment = PP_ALIGN.CENTER
    
    # Section title
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(3.5), Inches(11.333), Inches(1.2)
    )
    title_frame = title_box.text_frame
    title_frame.text = section_title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(48)
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS['white']
    title_para.alignment = PP_ALIGN.CENTER
    
    return slide

# ============================================================================
# SLIDE 1: Title Slide
# ============================================================================
add_title_slide(
    prs,
    "Passkey Authentication for CIAM-Based LCS Applications",
    "POC Demo • Identity Management Team • Senior Developer"
)

# ============================================================================
# SLIDE 2: Agenda
# ============================================================================
add_content_slide(prs, "Agenda", [
    "Today's Journey:",
    "• The Password Problem & Security Context",
    "• FIDO2 & Passkeys Fundamentals",
    "• How Passkeys Work (Registration & Authentication)",
    "• Security Benefits & Challenges",
    "• Our Current State: Azure AD B2C",
    "• The Gap: Why We Need Passkey Support Now",
    "• Solution: open-passkey Library",
    "• Implementation Architecture",
    "• Live Demo",
    "• Future Roadmap & Migration Path",
])

# ============================================================================
# SLIDE 3: Section Divider - The Problem
# ============================================================================
add_section_divider(prs, "The Password Problem", "01")

# ============================================================================
# SLIDE 4: Password Statistics
# ============================================================================
add_content_slide(prs, "Why Passwords Fail", [
    "The Alarming Statistics:",
    "• 77% of all hacking-related breaches involve stolen or reused credentials",
    "• 3,000% increase in AI-powered phishing attacks targeting corporate credentials",
    "• 48% of people have abandoned online purchases due to forgotten passwords",
    "• 65% of people reuse passwords across multiple sites",
    "• 36% have had at least one account compromised",
    "",
    "The Core Issue:",
    "• Passwords are shared secrets — server stores them (even if hashed)",
    "• Breach = credential theft at scale",
    "• Phishing works because users can't distinguish fake sites",
    "• No cryptographic binding to the service",
])

# ============================================================================
# SLIDE 5: Section Divider - FIDO2 Fundamentals
# ============================================================================
add_section_divider(prs, "FIDO2 & Passkeys Fundamentals", "02")

# ============================================================================
# SLIDE 6: What is FIDO2?
# ============================================================================
add_content_slide(prs, "What is FIDO2?", [
    "FIDO = Fast IDentity Online",
    "",
    "The FIDO Alliance:",
    "• Founded 2012 by PayPal, Lenovo, Google, Microsoft, Apple",
    "• 1000+ member organizations",
    "• Mission: Eliminate passwords through open standards",
    "",
    "FIDO2 = The Current Standard (2018)",
    "• WebAuthn (W3C): Browser-side JavaScript API",
    "• CTAP2: Client-to-Authenticator Protocol",
    "• Enables true passwordless authentication",
    "",
    "Key Principle:",
    "• Public-key cryptography instead of shared secrets",
    "• Private key NEVER leaves the device",
])

# ============================================================================
# SLIDE 7: What is a Passkey?
# ============================================================================
add_content_slide(prs, "What is a Passkey?", [
    "A passkey is a FIDO2 credential designed to replace passwords:",
    "",
    "🔑 Cryptographic Key Pair:",
    "• Private Key: Stored in Secure Enclave/TPM/TEE, NEVER leaves device",
    "• Public Key: Stored on server, useless for authentication",
    "",
    "✨ Key Properties:",
    "• Domain-bound to RP ID (e.g., 'github.com')",
    "• Discoverable credential (no username needed)",
    "• Phishing-resistant by design",
    "• Unique per service (no cross-site tracking)",
    "",
    "❌ What a Passkey is NOT:",
    "• Not a password to type",
    "• Not stored on server (only public key)",
    "• Not transferable to phishing sites",
    "• Not a shared secret",
])

# ============================================================================
# SLIDE 8: Section Divider - How It Works
# ============================================================================
add_section_divider(prs, "How Passkeys Work", "03")

# ============================================================================
# SLIDE 9: Registration Flow
# ============================================================================
add_content_slide(prs, "Registration Flow (One-Time Setup)", [
    "Step-by-Step Registration:",
    "",
    "1. User clicks 'Create Passkey'",
    "2. Server sends registration options (challenge, RP ID, user info)",
    "3. Browser calls navigator.credentials.create()",
    "4. Authenticator performs SECURE operation:",
    "   • Verify user identity (biometric/PIN)",
    "   • Generate new key pair",
    "   • Store private key securely",
    "   • Create Credential ID",
    "5. Authenticator returns credential (public key, attestation)",
    "6. Server validates and stores public key + credential ID",
    "7. Registration complete!",
    "",
    "Result: Private key stays on device, public key stored on server",
])

# ============================================================================
# SLIDE 10: Authentication Flow
# ============================================================================
add_content_slide(prs, "Authentication Flow (Every Login)", [
    "Step-by-Step Login:",
    "",
    "1. User clicks 'Sign in with Passkey'",
    "2. Server sends authentication options (FRESH challenge, RP ID)",
    "3. Browser calls navigator.credentials.get()",
    "4. Authenticator finds passkey for this domain",
    "5. User verifies locally (biometric/PIN)",
    "6. Authenticator signs the challenge with private key",
    "7. Returns assertion (credential ID, signature, authenticator data)",
    "8. Server verifies:",
    "   • Retrieve public key by credential ID",
    "   • Verify signature ✅",
    "   • Verify origin (domain binding) ✅",
    "   • Verify challenge (prevent replay) ✅",
    "9. Authentication successful! Issue JWT token",
])

# ============================================================================
# SLIDE 11: What Gets Stored Where
# ============================================================================
add_content_slide(prs, "What Gets Stored Where?", [
    "📱 User Device (e.g., iPhone, Laptop):",
    "• 🔒 Private Key (NEVER leaves device)",
    "   - Stored in Secure Enclave / TPM / TEE",
    "   - Used to SIGN challenges",
    "• 📋 Credential Metadata",
    "   - Credential ID",
    "   - RP ID (domain)",
    "   - User Handle",
    "   - Created timestamp",
    "",
    "🏢 Our Server:",
    "• 🔓 Public Key (SAFE to store - can't authenticate)",
    "   - Mathematical partner to private key",
    "   - Used to VERIFY signatures",
    "• 👤 User Record",
    "   - User ID",
    "   - Credential ID",
    "   - Public Key",
    "   - Sign counter",
])

# ============================================================================
# SLIDE 12: Section Divider - Security
# ============================================================================
add_section_divider(prs, "Security Benefits & Challenges", "04")

# ============================================================================
# SLIDE 13: Why Passkeys Are Unphishable
# ============================================================================
add_content_slide(prs, "Why Passkeys Are Unphishable", [
    "Attack Scenario → Passkey Defense:",
    "",
    "🎣 Phishing (fake github-login.com)",
    "   → Domain Binding: Passkey ONLY works for 'github.com'",
    "   → Fake site = automatic refusal. User protected automatically.",
    "",
    "💾 Database Breach (server hacked)",
    "   → Only PUBLIC key stored on server",
    "   → Public keys are USELESS for authentication",
    "   → Nothing to steal!",
    "",
    "🔄 Credential Stuffing (reused credentials)",
    "   → Each passkey is UNIQUE per site",
    "   → Stolen key from Site A cannot work on Site B",
    "",
    "👁️ Man-in-the-Middle (intercepted traffic)",
    "   → Challenge is signed locally",
    "   → Even if intercepted, signature only valid for THIS challenge",
    "",
    "🔁 Replay Attack (recorded login)",
    "   → Challenge is random and ONE-TIME",
    "   → Each login uses fresh challenge",
    "   → Captured responses are expired immediately",
])

# ============================================================================
# SLIDE 14: Security Comparison
# ============================================================================
add_content_slide(prs, "Security & UX Comparison", [
    "Authentication Method Comparison:",
    "",
    "🔑 Password Only:",
    "   Security: ⭐☆☆☆☆ | UX: ⭐⭐☆☆☆ | Phishing: ❌ NO",
    "",
    "🔑+📱 Password + SMS OTP:",
    "   Security: ⭐⭐⭐☆☆ | UX: ⭐⭐☆☆☆ | Phishing: ❌ NO (real-time relay)",
    "",
    "🔑+🔢 Password + TOTP (Authenticator App):",
    "   Security: ⭐⭐⭐⭐☆ | UX: ⭐⭐☆☆☆ | Phishing: ❌ NO (code can be phished)",
    "",
    "🔑 Passkey:",
    "   Security: ⭐⭐⭐⭐⭐ | UX: ⭐⭐⭐⭐⭐ | Phishing: ✅ YES",
    "",
    "Real-World Impact:",
    "• Amazon: 6x faster sign-in times",
    "• Google: 4x improvement in sign-in success rates",
    "• CVS Health: 98% reduction in mobile account takeover fraud",
    "• Average login: ~3-4 seconds (vs ~20s for passwords)",
])

# ============================================================================
# SLIDE 15: Challenges
# ============================================================================
add_content_slide(prs, "Challenges & Mitigations", [
    "⚠️ Challenges:",
    "",
    "1. Device Loss",
    "   → Mitigation: Backup methods (email, backup codes, hardware keys)",
    "",
    "2. Cross-Device Authentication",
    "   → Mitigation: BLE + QR code flow (WebAuthn standard)",
    "",
    "3. Algorithm Ecosystem",
    "   → ES256 today (broad compatibility)",
    "   → Monitor PQ/hybrid readiness for future",
    "",
    "4. Account Recovery",
    "   → Mitigation: Identity verification with support team",
    "",
    "5. Browser Compatibility",
    "   → Mitigation: Progressive enhancement (passkey primary, fallback options)",
    "",
    "6. Challenge Store Scalability",
    "   → Current: In-memory (POC)",
    "   → Production: Redis/DB for multi-instance deployments",
])

# ============================================================================
# SLIDE 16: Section Divider - Current State
# ============================================================================
add_section_divider(prs, "Our Current State: Azure AD B2C", "05")

# ============================================================================
# SLIDE 17: Current Architecture
# ============================================================================
add_content_slide(prs, "Current Authentication Architecture", [
    "Current Setup:",
    "",
    "• Identity Provider: Azure AD B2C (Custom Policies)",
    "• Token Issuance: B2C issues JWT tokens after authentication",
    "• Application: CIAM-based LCS applications",
    "• Flow: User → App → B2C → Token → App",
    "",
    "Current Limitations:",
    "",
    "❌ No Native Passkey Support",
    "   - B2C does not support FIDO2/WebAuthn",
    "   - Cannot offer modern passwordless authentication",
    "",
    "⏰ Retirement Timeline",
    "   - Azure AD B2C going out of service by 2030",
    "   - No new features being added",
    "",
    "🔮 Future Direction",
    "   - Planned migration to Microsoft Entra External ID",
    "   - Entra has native passkey support",
    "   - Migration timeline: 12-18 months",
])

# ============================================================================
# SLIDE 18: The Gap
# ============================================================================
add_content_slide(prs, "The Gap: Why We Need Passkey Support Now", [
    "The Problem:",
    "",
    "Current State:",
    "• Azure AD B2C: No passkey support",
    "• Timeline: B2C retiring by 2030",
    "• Future: Moving to Microsoft Entra (has native passkey support)",
    "• Problem: Entra migration is 12-18 months away",
    "",
    "We need passkey support NOW, not in 18 months!",
    "",
    "Why Wait is Not an Option:",
    "",
    "• User Experience: Passwords create friction → abandonment",
    "• Security Risk: Password-based auth vulnerable to modern attacks",
    "• Competitive Pressure: Industry moving to passkeys",
    "   - Apple, Google, Microsoft, Amazon, GitHub all support them",
    "• Compliance: NIST 800-63B recommends phishing-resistant auth",
    "• Future-Proofing: Need to prepare for B2C retirement",
])

# ============================================================================
# SLIDE 19: Section Divider - Solution
# ============================================================================
add_section_divider(prs, "Solution: Bridge Architecture", "06")

# ============================================================================
# SLIDE 20: Solution Approach
# ============================================================================
add_content_slide(prs, "Our Approach: Bridge Architecture", [
    "Bridge Architecture Concept:",
    "",
    "Implement passkey verification NOW with open-passkey,",
    "continue using current token authority (B2C) for downstream compatibility,",
    "keep migration flexibility for future Entra target state.",
    "",
    "✅ Immediate Value:",
    "• Deploy passkeys now, independent of Entra migration",
    "• Zero disruption to existing token flow (B2C still issues tokens)",
    "• Flexibility to migrate to Entra-native passkeys later",
    "• Cost-effective: Open-source vs commercial providers",
    "• Full control over UX and error handling",
    "",
    "🔄 Future Path:",
    "• Phase 1 (Now): open-passkey bridge (this POC)",
    "• Phase 2 (2026): Production hardening (Redis, monitoring, tests)",
    "• Phase 3 (2026+): Migrate to Microsoft Entra (native passkey support)",
    "• Phase 4 (2027+): Simplify architecture (remove bridge)",
])

# ============================================================================
# SLIDE 21: Section Divider - Library
# ============================================================================
add_section_divider(prs, "Library Choice: open-passkey", "07")

# ============================================================================
# SLIDE 22: Library Overview
# ============================================================================
add_content_slide(prs, "Why open-passkey?", [
    "Library: https://github.com/locke-inc/open-passkey",
    "License: MIT (permissive, commercial-friendly)",
    "Language: Java (Spring Boot integration available)",
    "Version: v0.1.4",
    "",
    "✅ Key Strengths:",
    "",
    "1. Production-Ready ES256",
    "   • ECDSA P-256 (secp256r1) + SHA-256",
    "   • Broad browser/authenticator support",
    "   • Reliable baseline for production",
    "",
    "2. Spring Boot Integration",
    "   • First-class Spring support",
    "   • Reduces boilerplate significantly",
    "",
    "3. Clean Abstractions",
    "   • CredentialStore and ChallengeStore interfaces",
    "   • Easy to implement custom stores (PostgreSQL, Redis)",
    "",
    "4. Active Development",
    "   • MIT license ensures freedom to modify",
    "   • Standards-compliant (WebAuthn spec)",
    "",
    "5. Flexible Architecture",
    "   • Can exclude auto-configuration",
    "   • Build custom endpoints",
    "   • Multi-RP support built-in",
])

# ============================================================================
# SLIDE 23: ES256 Algorithm Choice
# ============================================================================
add_content_slide(prs, "Why ES256 Algorithm?", [
    "ES256 = ECDSA with P-256 (secp256r1) curve + SHA-256",
    "",
    "Algorithm Comparison:",
    "",
    "✅ ES256 (Selected):",
    "   • Identifier: -7",
    "   • Broad browser/authenticator support",
    "   • Reliable baseline for production",
    "   • Lower ecosystem risk",
    "",
    "⚠️ RS256:",
    "   • Larger signatures, slower",
    "   • Less common in passkey ecosystem",
    "",
    "⚠️ EdDSA:",
    "   • Newer algorithm",
    "   • Not universally supported yet",
    "",
    "🔮 PQ/Hybrid (Future):",
    "   • Post-quantum direction present in library",
    "   • Gated by client/browser/authenticator maturity",
    "   • Server-side support exists, but ecosystem not ready",
    "",
    "Decision: ES256 is the practical choice for today's deployments.",
])

# ============================================================================
# SLIDE 24: Section Divider - Architecture
# ============================================================================
add_section_divider(prs, "Implementation Architecture", "08")

# ============================================================================
# SLIDE 25: High-Level Architecture
# ============================================================================
add_content_slide(prs, "High-Level Architecture", [
    "System Components:",
    "",
    "Frontend (localhost:5555):",
    "• Static UI (HTML/JS)",
    "• WebAuthn API calls (navigator.credentials.create/get)",
    "• Sends credential/assertion payloads to backend",
    "",
    "Backend (localhost:8080):",
    "• CustomPasskeyController: /passkey/** endpoints",
    "• PasskeyServiceRegistry: Multi-RP support",
    "• open-passkey layer: WebAuthn verification",
    "• PostgresCredentialStore: Persist credentials",
    "• TokenService: Handoff to B2C for JWT issuance",
    "",
    "Database:",
    "• PostgreSQL: passkey_credentials table",
    "• Stores: credential_id, user_id, public_key_cose, sign_count",
    "",
    "Token Flow:",
    "• Passkey auth succeeds → extract userId",
    "• Call B2C token endpoint (client_credentials + username)",
    "• Return JWT to frontend",
    "• Existing downstream APIs work unchanged",
])

# ============================================================================
# SLIDE 26: Key Components
# ============================================================================
add_content_slide(prs, "Key Components", [
    "1. CustomPasskeyController",
    "   • Endpoints: /passkey/register/begin, /passkey/register/finish",
    "   • Endpoints: /passkey/login/begin, /passkey/login/finish",
    "   • Orchestrates WebAuthn ceremonies",
    "   • Handles token handoff after verification",
    "   • Multi-RP routing by rpId parameter",
    "",
    "2. PasskeyServiceRegistry",
    "   • One PasskeyService per RP (Relying Party)",
    "   • Configuration: passkey.relying-parties in application.yaml",
    "   • Validates incoming rpId, rejects unknown with 400",
    "",
    "3. PostgresCredentialStore",
    "   • Persists credential metadata in PostgreSQL",
    "   • Schema: passkey_credentials table",
    "   • Operations: store, get, getByUser, update, delete",
    "",
    "4. TokenService (Identity Token Handoff)",
    "   • After passkey verification, fetch JWT from B2C",
    "   • Uses WebClient for HTTP call",
    "   • Returns token payload to frontend",
])

# ============================================================================
# SLIDE 27: Configuration
# ============================================================================
add_content_slide(prs, "Configuration Example", [
    "application.yaml:",
    "",
    "passkey:",
    "  relying-parties:",
    "    - rp-id: localhost",
    "      rp-display-name: My App (Local)",
    "      origin: http://localhost:5555",
    "      additional-origins:",
    "        - http://localhost:3000",
    "        - http://localhost:4444",
    "",
    "authorization-server:",
    "  token:",
    "    endpoint: https://LexmarkB2CDevelopment.b2clogin.com/...",
    "    client-id: ${CLIENT_ID}",
    "    client-secret: ${CLIENT_SECRET}",
    "    scope: https://.../.default",
    "",
    "Key Points:",
    "• Multi-origin support for local development",
    "• CORS origins derived from RP configuration",
    "• Token endpoint configurable per environment",
    "• Secrets injected via environment variables",
])

# ============================================================================
# SLIDE 28: Database Schema
# ============================================================================
add_content_slide(prs, "Database Schema", [
    "passkey_credentials table:",
    "",
    "CREATE TABLE passkey_credentials (",
    "    id              BIGSERIAL PRIMARY KEY,",
    "    credential_id   VARCHAR(512)  NOT NULL UNIQUE,",
    "    user_id         VARCHAR(256)  NOT NULL,",
    "    public_key_cose BYTEA         NOT NULL,",
    "    sign_count      BIGINT        NOT NULL DEFAULT 0,",
    "    prf_salt        BYTEA,",
    "    prf_supported   BOOLEAN       NOT NULL DEFAULT FALSE,",
    "    created_at      TIMESTAMPTZ   NOT NULL DEFAULT NOW(),",
    "    last_used_at    TIMESTAMPTZ",
    ");",
    "",
    "Indexes:",
    "• idx_pk_user_id ON passkey_credentials(user_id)",
    "• idx_pk_credential_id ON passkey_credentials(credential_id)",
    "",
    "Key Fields:",
    "• credential_id: Base64URL-encoded (matches WebAuthn wire format)",
    "• public_key_cose: COSE-format public key (bytea)",
    "• sign_count: Increments on each use (prevents cloning)",
    "• prf_salt: PRF extension support (future-proofing)",
])

# ============================================================================
# SLIDE 29: Section Divider - Demo
# ============================================================================
add_section_divider(prs, "Live Demo", "09")

# ============================================================================
# SLIDE 30: Demo Flow
# ============================================================================
add_content_slide(prs, "Demo Flow", [
    "What We'll Show:",
    "",
    "Scene 1: Login with Passkey (2 min)",
    "• Open browser to localhost:5555",
    "• Show login page with two options",
    "• Click 'Sign in with Passkey'",
    "• Browser biometric prompt",
    "• Success: Redirect to My Account page",
    "• Show JWT token in DevTools",
    "",
    "Scene 2: Registration Flow (3 min)",
    "• Navigate to registration page",
    "• Enter email address",
    "• Browser creates passkey",
    "• Biometric confirmation",
    "• Show database record in PostgreSQL",
    "",
    "Scene 3: Architecture Deep Dive (5 min)",
    "• Show backend logs (Spring Boot console)",
    "• Show database (PostgreSQL credentials)",
    "• Show token handoff (network call to B2C)",
    "• Explain multi-RP configuration",
    "",
    "Scene 4: Security Demonstration (3 min)",
    "• Show origin validation",
    "• Explain domain binding",
    "• Show sign counter increment",
])

# ============================================================================
# SLIDE 31: Demo Screenshots
# ============================================================================
add_content_slide(prs, "Demo: Login Page", [
    "Login Page Features:",
    "",
    "Two Authentication Options:",
    "",
    "1. Sign in with Passkey (Recommended)",
    "   • Green highlighted button",
    "   • Uses registered passkey for device",
    "   • No password needed",
    "   • Biometric verification",
    "",
    "2. Sign in with Password (Legacy Fallback)",
    "   • Redirects to B2C sign-in page",
    "   • Traditional username/password",
    "   • Maintains backward compatibility",
    "",
    "Email Field:",
    "• Optional for passkey login (discoverable credential)",
    "• Required for password fallback",
    "• 'Leave blank to use a discoverable passkey'",
    "",
    "User Experience:",
    "• Clean, modern UI",
    "• Clear visual hierarchy",
    "• Fast 3-4 second login with passkey",
])

# ============================================================================
# SLIDE 32: My Account Page
# ============================================================================
add_content_slide(prs, "Demo: My Account Page", [
    "After Successful Authentication:",
    "",
    "Personal Information:",
    "• Display Name: Sandeep OA-CIAM-US",
    "• E-mail: sandeep-ciam-oa-us@test.onekx.co",
    "• First Name: sandeep",
    "• Last Name: mohanty",
    "• Department: B2C Department",
    "• Cost Center: B2C Cost Center",
    "",
    "Actions Available:",
    "• Reset Password",
    "• Register Passkey (add another device)",
    "",
    "Assigned Groups:",
    "• Admin",
    "• test PAs",
    "",
    "Printer Login:",
    "• Login Method: PIN, Badges or Secure Login",
    "• Badges management",
    "• PIN reset",
    "",
    "Key Point:",
    "• User is authenticated with passkey",
    "• JWT token from B2C is used for API calls",
    "• All existing functionality works unchanged",
])

# ============================================================================
# SLIDE 33: Section Divider - Roadmap
# ============================================================================
add_section_divider(prs, "Future Roadmap", "10")

# ============================================================================
# SLIDE 34: Migration Timeline
# ============================================================================
add_content_slide(prs, "Migration Path to Microsoft Entra", [
    "Phased Migration Strategy:",
    "",
    "Phase 1: Current (Q1-Q2 2025)",
    "• open-passkey Bridge Architecture (This POC)",
    "• Deploy passkey support immediately",
    "• Independent of Entra migration timeline",
    "",
    "Phase 2: Production Hardening (Q3-Q4 2025)",
    "• Replace in-memory challenge store with Redis",
    "• Add comprehensive controller + integration tests",
    "• Implement timeout/retry/circuit-breaker for token handoff",
    "• Add observability (metrics, alerts, correlation IDs)",
    "• Enable HTTPS-only origins",
    "• Implement account recovery flows",
    "• Security audit by IAM team",
    "",
    "Phase 3: Entra Migration (2026)",
    "• Migrate to Microsoft Entra External ID",
    "• Leverage native passkey support",
    "• Simplify architecture (remove bridge)",
    "",
    "Phase 4: Optimization (2027+)",
    "• Entra-native flow",
    "• Remove open-passkey dependency",
    "• Full Microsoft ecosystem integration",
])

# ============================================================================
# SLIDE 35: Production Checklist
# ============================================================================
add_content_slide(prs, "Production Checklist", [
    "Before Production Deployment:",
    "",
    "Identity & Security:",
    "☐ Validate passkey → token handoff policy with IAM/security stakeholders",
    "☐ Ensure RP ID/origin exactness for each environment",
    "☐ Use HTTPS-only origins for production",
    "☐ Keep secrets only in vault/secret manager",
    "☐ Implement account recovery and device-loss flows",
    "",
    "Architecture & Scale:",
    "☐ Replace/augment MemoryChallengeStore for multi-instance deployments",
    "☐ Add timeout/retry/circuit-breaker around token handoff",
    "☐ Define migration roadmap from custom bridge to Entra-native flow",
    "",
    "Reliability & Testing:",
    "☐ Add controller and integration tests for all critical paths",
    "☐ Add browser/device compatibility matrix tests",
    "☐ Add regression checks when updating open-passkey JARs",
    "",
    "Observability:",
    "☐ Emit metrics for begin/finish pass/fail by RP",
    "☐ Alert on spikes in authentication/token failures",
    "☐ Add correlation IDs across passkey and token calls",
    "",
    "Dependency Governance:",
    "☐ Track open-passkey source tag/commit used for each bundled JAR",
    "☐ Establish regular upgrade review cadence",
    "☐ Monitor CVEs for crypto/CBOR/web stack dependencies",
])

# ============================================================================
# SLIDE 36: Key Takeaways
# ============================================================================
add_content_slide(prs, "Key Takeaways", [
    "What We've Achieved:",
    "",
    "🎯 Immediate Value:",
    "• Deploy passkeys NOW, independent of Entra timeline",
    "• Zero disruption to existing B2C token flow",
    "• Cost-effective open-source solution",
    "",
    "🛡️ Security Benefits:",
    "• 100% phishing resistant (domain-bound credentials)",
    "• Server breach safe (only public keys stored)",
    "• No passwords to steal or phish",
    "• 6x faster login, 4x better success rate",
    "",
    "🏗️ Technical Approach:",
    "• Bridge architecture with open-passkey library",
    "• Token handoff to B2C (existing APIs work unchanged)",
    "• Multi-RP support for flexible deployment",
    "• Clear migration path to Microsoft Entra",
    "",
    "🚀 Future Ready:",
    "• Standards-based (FIDO2/WebAuthn)",
    "• Industry alignment (Apple, Google, Microsoft, Amazon)",
    "• Prepared for B2C retirement (2030)",
])

# ============================================================================
# SLIDE 37: Business Case
# ============================================================================
add_content_slide(prs, "The Business Case", [
    "Investment & Impact:",
    "",
    "Security:",
    "• Eliminate password-related breaches (77% of hacking incidents)",
    "• Meet NIST 800-63B (phishing-resistant authentication)",
    "• Reduce account takeover fraud (CVS: 98% reduction)",
    "",
    "User Experience:",
    "• 6x faster login (3-4 seconds vs 20+ seconds)",
    "• 4x better sign-in success rate (98% vs 70%)",
    "• 48% reduction in abandonment due to forgotten passwords",
    "• No more password reset support tickets",
    "",
    "Cost:",
    "• Reduce password reset support costs",
    "• Lower helpdesk burden",
    "• Open-source library vs commercial providers",
    "",
    "Compliance:",
    "• Align with modern authentication standards",
    "• Prepare for B2C retirement (2030)",
    "• Ready for Entra migration",
    "",
    "Competitive:",
    "• Match industry leaders (Amazon, Google, GitHub)",
    "• Modern, innovative identity solution",
])

# ============================================================================
# SLIDE 38: Q&A
# ============================================================================
add_title_slide(
    prs,
    "Questions & Answers",
    "POC Demo • Passkey Authentication for CIAM"
)

# ============================================================================
# SLIDE 39: Contact & Resources
# ============================================================================
add_content_slide(prs, "Contact & Resources", [
    "Presenter: Sandeep Mohanty",
    "Role: Senior Developer, Identity Management Team",
    "Email: sandeep-ciam-oa-us@test.onekx.co",
    "",
    "Resources:",
    "",
    "• open-passkey Library:",
    "  https://github.com/locke-inc/open-passkey",
    "",
    "• FIDO Alliance:",
    "  https://fidoalliance.org/",
    "",
    "• WebAuthn Specification:",
    "  https://www.w3.org/TR/webauthn-2/",
    "",
    "• Passkeys.com:",
    "  https://passkeys.com/",
    "",
    "• Microsoft Entra External ID:",
    "  https://learn.microsoft.com/en-us/entra/external-id/",
    "",
    "Demo Environment:",
    "• Backend: localhost:8080 (Spring Boot + open-passkey)",
    "• Frontend: localhost:5555 (Static HTML/JS)",
    "• Database: PostgreSQL (passkey_credentials table)",
])

# ============================================================================
# SLIDE 40: Thank You
# ============================================================================
add_title_slide(
    prs,
    "Thank You!",
    "Passkey Authentication POC • Identity Management Team"
)

# ============================================================================
# Save presentation
# ============================================================================
output_path = "d:/knowledge-base/Identity/FIDO2/Passkey_Authentication_POC_Demo.pptx"
prs.save(output_path)
print(f"✅ Presentation saved to: {output_path}")
print(f"   Total slides: {len(prs.slides)}")